"""PDF to PowerPoint converter processor."""

import asyncio
from pathlib import Path

from app.processors.base import BaseProcessor
from app.core.exceptions import ProcessingError


class PDFToPowerPointConverter(BaseProcessor):
    """Convert PDF to PowerPoint presentation."""

    async def execute(
        self,
        input_path: Path,
        output_path: Path,
        dpi: int = 150,
    ) -> Path:
        """
        Convert PDF to PowerPoint.

        Each PDF page becomes a slide with the page rendered as an image.

        Args:
            input_path: Path to the input PDF
            output_path: Path for the output PPTX
            dpi: Resolution for rendering pages (higher = better quality, larger file)

        Returns:
            Path to the PowerPoint file
        """

        def _convert():
            from pdf2image import convert_from_path
            from pptx import Presentation
            from pptx.util import Inches, Pt
            import tempfile
            import os

            # Convert PDF pages to images
            images = convert_from_path(str(input_path), dpi=dpi)

            if not images:
                raise ProcessingError("No pages found in PDF")

            # Create PowerPoint presentation
            prs = Presentation()

            # Set slide dimensions based on first image aspect ratio
            first_img = images[0]
            img_width, img_height = first_img.size
            aspect_ratio = img_width / img_height

            # Standard slide size (10x7.5 inches)
            slide_width = Inches(10)
            slide_height = Inches(7.5)

            prs.slide_width = slide_width
            prs.slide_height = slide_height

            # Create temp directory for images
            with tempfile.TemporaryDirectory() as temp_dir:
                for idx, image in enumerate(images):
                    # Save image temporarily
                    img_path = os.path.join(temp_dir, f"page_{idx}.png")
                    image.save(img_path, "PNG")

                    # Add blank slide
                    blank_layout = prs.slide_layouts[6]  # Blank layout
                    slide = prs.slides.add_slide(blank_layout)

                    # Calculate image size to fit slide while maintaining aspect ratio
                    img_aspect = image.width / image.height
                    slide_aspect = slide_width / slide_height

                    if img_aspect > slide_aspect:
                        # Image is wider - fit to width
                        width = slide_width
                        height = int(slide_width / img_aspect)
                    else:
                        # Image is taller - fit to height
                        height = slide_height
                        width = int(slide_height * img_aspect)

                    # Center the image
                    left = int((slide_width - width) / 2)
                    top = int((slide_height - height) / 2)

                    # Add image to slide
                    slide.shapes.add_picture(img_path, left, top, width, height)

            # Save presentation
            output_path_pptx = output_path.with_suffix(".pptx")
            prs.save(str(output_path_pptx))

            return output_path_pptx

        try:
            return await asyncio.to_thread(_convert)
        except ImportError as e:
            raise ProcessingError(
                f"Required package not installed: {e}. "
                "Install with: pip install pdf2image python-pptx"
            )
        except Exception as e:
            raise ProcessingError(f"Failed to convert PDF to PowerPoint: {str(e)}")
